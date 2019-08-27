'''
createPPTX creates a powerpoint from a given sample powerpoint.
The source will come from a project approval.
'''

# imports go here
from pptx import Presentation
# from pptx.util import Inches
import argparse
import pandas as pd
from datetime import date

# Classes and globals

# Placeholder layout and indices for the Title Slide
TITLE_SLIDE_LAYOUT = 0
TITLE_SLIDE = {'subtitle': 1}

BLANK_SLIDE_LAYOUT = 7

# Placeholder indices for the first approval slide
APP_SLIDE1_LAYOUT = 2
APP_SLIDE1 = {'Key': 27,
              # 'ID': #,
              # 'Created': #,
              # 'Title': #,
              'Department': 23,
              # 'Classification' #,
              'Description': 1,
              'Status': 24,
              'Return Type': 25,
              'ROI': 15,
              # 'Break Even Period': #,
              'Return': 16,
              # 'Duration': #,
              'Total Investment': 17,
              # 'Labor T-Shirt Size': #,
              # 'Labor Investment': #,
              # 'Non-Labor Investment': #,
              'Sponsor': 19,
              'Project Owner': 18,
              'Line of Business': 21,
              # 'States': 22,
              }

# Placeholder indices for the second approval slide
APP_SLIDE2_LAYOUT = 3
APP_SLIDE2 = {'Return Description': 1,
              'Scope': 14,
              'Scope Exclusions': 15,
              'Systems': 17,
              'Business Area Impacts': 16,
              'Investment Description': 13,
              }

# Data Column Names for simplified mapping
COLMUN_NAMES = {'Issue key': 'Key',
                'Issue id': 'ID',
                'Created': 'Created',
                'Summary': 'Title',
                'Custom field (Sponsor Department)': 'Department',
                'Custom field (Project Classification)':
                    'Classification',
                'Description': 'Description',
                'Status': 'Status',
                'Custom field (Return Type)': 'Return Type',
                'Custom field (Return on Investment)': 'ROI',
                'Custom field (Break Even Period)': 'Break Even Period',
                'Custom field (Return)': 'Return',
                'Custom field (Duration)': 'Duration',
                'Custom field (Total Investment)': 'Total Investment',
                'Custom field (Labor T-Shirt Size)': 'Labor T-Shirt Size',
                'Custom field (Labor Investment)': 'Labor Investment',
                'Custom field (Non-Labor Investment)': 'Non-Labor Investment',
                'Custom field (Executive Sponsor)': 'Sponsor',
                'Custom field (Sponsor)': 'Project Owner',
                'Custom field (Business Area Impacts)':
                    'Business Area Impacts',
                'Custom field (Labor Investment Description)':
                    'Investment Description',
                'Custom field (Line of Business)': 'Line of Business',
                'Custom field (Return Description)': 'Return Description',
                'Custom field (Scope)': 'Scope',
                'Custom field (Scope Exclusions)': 'Scope Exclusions',
                'Custom field (Systems)': 'Systems',
                }


# functions go here
def create_ppt_title_slide(prs, title_txt, subtitle_txt):
    '''Creates a title slide in the presentation'''
    # choose the Title layout for the slide
    slide_layout = prs.slide_layouts[TITLE_SLIDE_LAYOUT]
    slide = prs.slides.add_slide(slide_layout)
    # assign title text
    title = slide.shapes.title
    title.text = title_txt
    # assign subtitle text
    subtitle = slide.placeholders[TITLE_SLIDE['subtitle']]
    subtitle.text = subtitle_txt
    # Insert a blank slide
    slide_layout = prs.slide_layouts[BLANK_SLIDE_LAYOUT]
    slide = prs.slides.add_slide(slide_layout)


def placeTextInSlide(slide, placeholderIndex, text):
    '''Inserts text into a slide.
    Handles exceptions with blanks.'''
    try:
        slide.placeholders[placeholderIndex].text = text
    except TypeError:
        slide.placeholders[placeholderIndex].text = 'Not Available'


def create_ppt_approval_slide(prs, title_txt, data):
    '''Creates a project approval slide in the presentation'''

    # choose the Title layout for the slide
    slide_layout = prs.slide_layouts[APP_SLIDE1_LAYOUT]
    slide = prs.slides.add_slide(slide_layout)

    # assign title text
    slide.shapes.title.text = data["Title"]

    # assign data to placeholder text from the sample ppt
    for key, value in APP_SLIDE1.items():
        placeTextInSlide(slide, value, data[key])

    slide_layout = prs.slide_layouts[APP_SLIDE2_LAYOUT]
    slide = prs.slides.add_slide(slide_layout)

    # assign title text
    slide.shapes.title.text = data["Title"]

    # assign data to placeholder text from the sample ppt
    for key, value in APP_SLIDE2.items():
        placeTextInSlide(slide, value, data[key])


def getDataFrame(data_file):
    data_frame = pd.read_csv(data_file)
    data_frame.rename(columns=COLMUN_NAMES, inplace=True)
    return data_frame


def create_pptx(data_file, input_file, output_file):
    '''Take the input powerpoint file and use it as the template
    for the output file.'''

    prs = Presentation(input_file)

    # Titles slide
    title_txt = "Project Approvals"
    subtitle_txt = "Created on {:%m-%d-%Y}".format(date.today())
    print('Creating title slide.')
    create_ppt_title_slide(prs, title_txt, subtitle_txt)

    # get the data from the data file
    print('Creating data frame.')
    df = getDataFrame(data_file)

    # Project Approvals Slides
    for index, row in df.iterrows():
        print('Creating approval slide ', index + 1, ' of ',
              len(df.index), '.')
        create_ppt_approval_slide(prs, "Project Approvals", row)

    prs.save(output_file)


if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("dataSource", type=str,
                        help="data source")
    parser.add_argument("pptxSource", type=str,
                        help="source pptx")
    parser.add_argument("pptxDestination", type=str,
                        help="destination pptx")
    args = parser.parse_args()
    create_pptx(args.dataSource, args.pptxSource, args.pptxDestination)
