'''
createPPTX creates a powerpoint from a given sample powerpoint.
The source will come from a project approval.
'''

# imports go here
from pptx import Presentation
# from pptx.util import Inches
import argparse
import pandas as pd
import numpy as np
import os
from datetime import date

# Classes and globals
DOUBLE_SIDED_PRINTING = False

PPT_TEMPLATE = "project_request_template.pptx"

# Placeholder layout and indices for the Title Slide
TITLE_SLIDE_LAYOUT = 0
TITLE_SLIDE = {'subtitle': 1}

# Placeholder layout for blank slides
BLANK_SLIDE_LAYOUT = 4

PAP_ALL_SLIDES_MAPPING = {'Summary': 'Title'}

# Placeholder indices for the first approval slide
PAP_SLIDE1_LAYOUT = 1
PAP_SLIDE1_COLUMN_MAPPING = {'Key': 27,
                             'Sponsor Department': 23,
                             'Description': 1,
                             'Status': 24,
                             'Return Type': 25,
                             'ROI': 15,
                             'Return': 16,
                             'Total Investment': 17,
                             'Sponsor': 18,
                             'Line of Business': 21,
                             'Return Description': 28,
                             'Investment Description': 13,
                             'Project Manager': 20,
                             'States': 22,
                             }

# Placeholder indices for the second approval slide
PAP_SLIDE2_LAYOUT = 2
PAP_SLIDE2_COLUMN_MAPPING = {'Key': 27,
                             'Scope': 14,
                             'Scope Exclusions': 15,
                             'Systems': 17,
                             'Business Area Impacts': 16,
                             }

REQUEST_SLIDE_LAYOUT = 3
REQUEST_SLIDE_COLUMN_MAPPING = {'Key': 27,
                                'Sponsor Department': 23,
                                'Description': 1,
                                'Status': 24,
                                'Sponsor': 18,
                                'Line of Business': 21,
                                'Return Description': 28,
                                'Project Manager': 20,
                                'States': 22,
                                'Return Type': [29, 30]
                                }

# Data Column Names for simplified mapping
COLUMNS = {'Issue key': 'Key',
           'Issue id': 'ID',
           'Created': 'Created',
           'Summary': 'Summary',
           'Custom field (Sponsor Department)': 'Sponsor Department',
           'Custom field (Project Classification)':
               'Classification',
           'Description': 'Description',
           'Status': 'Status',
           'Custom field (Return Type)': 'Return Type',
           'Custom field (Return on Investment)': 'ROI',
           'Custom field (Break Even Period)': 'Break Even Period',
           'Custom field (Return)': 'Return',
           'Custom field (Duration)': 'Duration',
           'Custom field (Total Investment)': 'Total Investment',
           'Custom field (Labor T-Shirt Size)': 'Labor T-Shirt Size',
           'Custom field (Labor Investment)': 'Labor Investment',
           'Custom field (Non-Labor Investment)': 'Non-Labor Investment',
           'Custom field (Sponsor)': 'Sponsor',
           'Custom field (Business Area Impacts)': 'Business Area Impacts',
           'Custom field (Labor Investment Description)':
               'Investment Description',
           'Custom field (Line of Business)': 'Line of Business',
           'Custom field (Return Description)': 'Return Description',
           'Custom field (Scope)': 'Scope',
           'Custom field (Scope Exclusions)': 'Scope Exclusions',
           'Custom field (Systems)': 'Systems',
           'Custom field (Project Manager)': 'Project Manager',
           }

STATUS_MAPPING = {"Project Request Approval": ["New Request",
                                               "More Details Required",
                                               "Sponsor Review",
                                               "Executive Sponsor Review",
                                               "Executive Committee Review",
                                               "Project Approved",
                                               "Future Consideration"],
                  "Project Schedule Approval": ["Scope Project",
                                                "Estimate Cost and Duration",
                                                "ROI Validation",
                                                "Portfolio Scheduling"],
                  "Inflight Projects": ["Project Scheduled",
                                        "Project In Progress"],
                  "On Hold Projects": ["Project On Hold"],
                  "Complete Projects": ["Project Complete"],
                  "Rejected Projects": ["Operational", "Project Rejected"],
                  }

COLUMNS_WITH_NAMES = ['Project Manager', 'Sponsor']


def createSlide(prs, layout_index):
    '''returns a new slide based on the slide index'''
    # add the slide from based on the layout
    slide_layout = prs.slide_layouts[layout_index]
    return prs.slides.add_slide(slide_layout)


def populateSlideFromList(slide, title_txt, placeholder_list):
    '''Populates a slide in the presentation based on a title,
    and a list of placeholder index and text pairs.
    Empty text will skip insertion of that element.'''
    # insert the title into the slide
    if title_txt:
        slide.shapes.title.text = title_txt
    # place the values into the slide
    for placeholder_index, txt in placeholder_list:
        placeTextInSlide(slide, placeholder_index, txt)


def populateSlideFromSeries(slide, title_txt, placeholder_map, series):
    '''Populates a slide based on the map of placeholders to a dataframe series.
    Ensures the data_frame contains the necessary columns. Unknown columns are
    removed.'''
    # insert the title into the slide
    if title_txt:
        slide.shapes.title.text = title_txt
    # reduce the placeholder map down to the columns that exsist in the series
    reduced_placeholder_map = {column_name: placeholder_map[column_name]
                               for column_name in placeholder_map.keys()
                               if column_name in series.keys()}
    # map the data columns to the placeholders
    for column_name, placeholder_index in reduced_placeholder_map.items():
        placeTextInSlide(slide, placeholder_index, series[column_name])


def placeParagraphInPlaceholder(slide, placeholder_index, text):
    ''' Places a paragraph in the placeholder list or integer
    and puts the text in the paragraph'''
    if type(placeholder_index) is list:
        for index in placeholder_index:
            paragraph = slide.placeholders[index].text_frame.add_paragraph()
            paragraph.text = text
    elif type(placeholder_index) is int:
        paragraph = slide.placeholders[placeholder_index] \
            .text_frame.add_paragraph()
        paragraph.text = text
    else:
        print("Placeholder", placeholder_index, "doesn't exist.")


def placeTextInPlaceholder(slide, placeholder_index, text):
    ''' Places text in a placeholder list or integer.'''
    if type(placeholder_index) is list:
        for index in placeholder_index:
            slide.placeholders[index].text_frame.text = text
    elif type(placeholder_index) is int:
        slide.placeholders[placeholder_index].text_frame.text = text
    else:
        print("Placeholder", placeholder_index, "doesn't exist.")


def placeTextInSlide(slide, placeholder_index, text):
    '''Inserts text into a slide.
    Handles exceptions with blanks.'''
    try:
        text_lines = str(text).splitlines()
        line_count = 0
        for line in text_lines:
            if line_count == 0:
                if line == 'nan':
                    placeTextInPlaceholder(slide, placeholder_index, '')
                else:
                    placeTextInPlaceholder(slide, placeholder_index, line)
            elif line != '':
                placeParagraphInPlaceholder(slide, placeholder_index, line)
            line_count += 1
    except TypeError:
        placeTextInPlaceholder(slide, placeholder_index, 'Not Available')
    except KeyError:
        return


def create_title_slides(presentation, title_txt, subtitle_txt):
    ''' Create a title slide '''
    # create the title slide in the presentation
    placeholder_list = [(TITLE_SLIDE['subtitle'], subtitle_txt)]
    # populate the placeholders in the title slide with the subtitle
    populateSlideFromList(createSlide(presentation, TITLE_SLIDE_LAYOUT),
                          title_txt, placeholder_list)
    # create a blank slide in between slides to support double sided printing
    if DOUBLE_SIDED_PRINTING:
        createSlide(presentation, BLANK_SLIDE_LAYOUT)


def cleanReturnTypeColumnForProjectRequests(data):
    '''Removes all values from the return type column except Compliance'''
    data['Return Type'] = (lambda x: x if x == 'Compliance' else ' ')(
                            data['Return Type'])
    return data


def create_project_request_slides(presentation, title_txt, data):
    ''' Creates a project request slide pair. '''
    # create the project request first slide in the presentation
    project_request_slide = createSlide(presentation, REQUEST_SLIDE_LAYOUT)
    # !modify dataframe to support fields needed for project requests
    data = cleanReturnTypeColumnForProjectRequests(data)
    # populate the placeholders in the project status slide with the data
    populateSlideFromSeries(project_request_slide, title_txt,
                            REQUEST_SLIDE_COLUMN_MAPPING,
                            data)
    # create a blank slide in between slides to support double sided printing
    if DOUBLE_SIDED_PRINTING:
        createSlide(presentation, BLANK_SLIDE_LAYOUT)


def create_project_status_slides(presentation, title_txt, data):
    ''' Create a project status slide pair.
    --- maybe pass in a list of the layouts '''
    # create the project status first slide in the presentation
    first_status_slide = createSlide(presentation, PAP_SLIDE1_LAYOUT)
    # populate the placeholders in the project status slide with the data
    populateSlideFromSeries(first_status_slide, title_txt,
                            PAP_SLIDE1_COLUMN_MAPPING,
                            data)
    # create the project status second slide in the presentation
    second_status_slide = createSlide(presentation, PAP_SLIDE2_LAYOUT)
    # populate the placeholders in the project status slide
    populateSlideFromSeries(second_status_slide, title_txt,
                            PAP_SLIDE2_COLUMN_MAPPING,
                            data)


def unique_list(list_of_lists):
    ''' Takes in a list of lists, combines them and returns the unique items'''
    # combine the lists
    combined_list = [item for sublist in list_of_lists for item in sublist]
    # use numpy to return the unique items in the list
    np_list = np.array(combined_list)
    return np.unique(np_list)


def convert_usernames_to_fullnames(data_frame):
    ''' compares a set of columnns against known usernames and converts them
        to full names '''
    # get the dictionary from the name_mapping.csv
    name_df = pd.read_csv("name_mapping.csv")
    name_dict = {name_df['Username'][ind]: name_df['Full Name'][ind]
                 for ind in name_df.index}
    # map the full names to the user names
    for column in COLUMNS_WITH_NAMES:
        data_frame[column] = data_frame[column].map(name_dict)
    return data_frame


def reduce_data_frame_columns(data_frame):
    ''' reduces the columns to just those needed for the pap'''
    # get only the needed column names for the slide decks
    needed_columns = unique_list([list(PAP_SLIDE1_COLUMN_MAPPING.keys()),
                                  list(PAP_SLIDE2_COLUMN_MAPPING.keys()),
                                  list(PAP_ALL_SLIDES_MAPPING.keys())])
    # filter the data_frame to just those needed columnns
    data_frame = data_frame.filter(needed_columns)
    return data_frame


def clean_data_frame(data_frame):
    ''' Cleans the data in the data_frame. '''
    # Update column names in the data_frame
    data_frame = data_frame.rename(columns=COLUMNS)
    # reduce the columns to those needed
    data_frame = reduce_data_frame_columns(data_frame)
    # change usernames to Full Names
    data_frame = convert_usernames_to_fullnames(data_frame)
    return data_frame


def getDataFrame(data_file):
    data_frame = pd.read_csv(data_file)
    data_frame = clean_data_frame(data_frame)
    return data_frame


def create_pap_pptx(df, prs, item, output_filename, sort_column):
    # Insert title slide with department and status_category such as
    #    project request, inflight or complete.

    for status_category in STATUS_MAPPING:
        statuses = STATUS_MAPPING[status_category]
        status_dataframe_subset = df.loc[df['Status'].isin(statuses)]
        if not status_dataframe_subset.empty:
            # Create section title slide
            title_txt = item + '\n' + status_category
            subtitle_txt = ('Created: {:%m-%d-%Y}'.format(date.today()))
            create_title_slides(prs, title_txt, subtitle_txt)

            # create the slides in status order
            for status in STATUS_MAPPING[status_category]:
                # subset the dataframe to the department rows
                df_subset = df.loc[(df['Status'] == status)]
                df_subset = df_subset.sort_values(by=["Summary"])
                df_subset = df_subset.sort_values(by=[sort_column])
                i = 0
                for index, series in df_subset.iterrows():
                    if (i == 0):
                        print(' *', status, '-', len(df_subset.index),
                              'slides')
                        # Proposed Project Slides
                    i += 1
                    print('  + Creating slide', i, 'of',
                          len(df_subset.index), '.')
                    title_txt = series["Summary"]
                    if status_category == "Project Request Approval":
                        create_project_request_slides(prs, title_txt, series)
                    else:
                        create_project_status_slides(prs, title_txt, series)

        prs.save(output_filename)


def create_pap_pptxs(data_frame, ppt_template, file_prefix, ppt_structure):
    '''Create a project approval ppt from a dataframe and ppt template.
        A pap ppt deck will be created for each dept with the output prefix.'''

    # make a folder for the ppts
    ppt_folder_path = file_prefix.replace('.pptx', '')
    if not os.path.exists(ppt_folder_path):
        os.makedirs(ppt_folder_path)

    # create a base path and filename for the ppts
    ppt_base_file_path = os.path.join(
                        ppt_folder_path,
                        os.path.basename(file_prefix))

    # slice determines the breakdown of project requests into ppts.
    slice = ppt_structure['Slice PPTs']
    sort = ppt_structure['Sort By']

    # any slice that is blank or nan should be marked unknown
    data_frame[slice] = data_frame[slice].replace(np.nan, 'Unknown',
                                                  regex=True)

    # create a ppt deck for each department
    for item in data_frame[slice].unique():
        # build a ppt deck for each of the departments
        ppt_deck = Presentation(ppt_template)

        # create the output filename using the file_prefix and department
        output_ppt_filename = ppt_base_file_path.replace(
                                '.pptx', '_' + str(item).replace(' ', '_') +
                                '.pptx')
        # get the department specific data
        item_dataframe = data_frame.loc[data_frame[slice] == item]

        # create the powerpoint for the department
        print('Creating', item, 'ppt.', len(item_dataframe), 'slides.')
        create_pap_pptx(item_dataframe, ppt_deck, item, output_ppt_filename,
                        sort)


if __name__ == "__main__":
    ''' This script pulls two arguments, one for the csv data source and
        the other for the output data file base directory and path'''

    # Pull two arguments - one for the source data and the other for the
    # output file folder name and prefix
    parser = argparse.ArgumentParser()
    parser.add_argument("dataSource", type=str,
                        help="data source")
    parser.add_argument("pptxDestination", type=str,
                        help="destination pptx")
    parser.add_argument('--dbl', dest='double', action='store_const',
                        const=True, default=False,
                        help='forces double sided printing (default: False)')
    args = parser.parse_args()

    # get the data from the data file
    print('Creating data frame.')
    pap_data_frame = getDataFrame(args.dataSource)

    # set double sided printing flag
    DOUBLE_SIDED_PRINTING = args.double

    # --- pass this in as an argument
    ppt_structure = {'Slice PPTs': 'Status',
                     'Sort By': 'Sponsor Department'}

    # print the PAP ppt
    print('Creating pap pptx.')
    create_pap_pptxs(pap_data_frame, PPT_TEMPLATE, args.pptxDestination,
                     ppt_structure)
